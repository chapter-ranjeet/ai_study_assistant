from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BASE_DIR = Path(r"C:\Users\Ranjeet Yadav\OneDrive\Desktop\ai-study-assistant")
OUTPUT_FILE = BASE_DIR / "AI_Study_Assistant_Project_Presentation.pptx"


def apply_title_style(run):
    run.font.name = "Calibri"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)


def apply_subtitle_style(run):
    run.font.name = "Calibri"
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = RGBColor(230, 230, 230)


def apply_heading_style(run):
    run.font.name = "Calibri"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)


def apply_body_style(run):
    run.font.name = "Calibri"
    run.font.size = Pt(22)
    run.font.color.rgb = RGBColor(240, 240, 240)


def set_dark_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(16, 24, 40)


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_dark_bg(slide)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "AI Study Assistant\n(Mentora AI)"
    for p in title.text_frame.paragraphs:
        for r in p.runs:
            apply_title_style(r)

    subtitle.text = (
        "Name : Ranjeet Yadav\n"
        "Course :- B-Tech Data Science\n"
        "Q-ID :- 25030861\n"
        "Year: 2nd\n"
        "Semester :- 4th"
    )
    for p in subtitle.text_frame.paragraphs:
        for r in p.runs:
            apply_subtitle_style(r)


def add_bulleted_slide(prs, title_text, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_dark_bg(slide)

    title = slide.shapes.title
    title.text = title_text
    for p in title.text_frame.paragraphs:
        for r in p.runs:
            apply_heading_style(r)

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    for i, bullet in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        for r in p.runs:
            apply_body_style(r)


def add_two_content_slide(prs, title_text, left_points, right_points):
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    set_dark_bg(slide)

    title = slide.shapes.title
    title.text = title_text
    for p in title.text_frame.paragraphs:
        for r in p.runs:
            apply_heading_style(r)

    left = slide.shapes.placeholders[1].text_frame
    right = slide.shapes.placeholders[2].text_frame
    left.clear()
    right.clear()

    for i, point in enumerate(left_points):
        p = left.paragraphs[0] if i == 0 else left.add_paragraph()
        p.text = point
        p.level = 0
        for r in p.runs:
            apply_body_style(r)

    for i, point in enumerate(right_points):
        p = right.paragraphs[0] if i == 0 else right.add_paragraph()
        p.text = point
        p.level = 0
        for r in p.runs:
            apply_body_style(r)


def add_image_slide(prs, title_text, image_path, caption):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_dark_bg(slide)

    title = slide.shapes.title
    title.text = title_text
    for p in title.text_frame.paragraphs:
        for r in p.runs:
            apply_heading_style(r)

    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(1.0), Inches(1.4), width=Inches(11.3), height=Inches(4.8))

    box = slide.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(12.0), Inches(0.8))
    tf = box.text_frame
    tf.text = caption
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = "Calibri"
            r.font.size = Pt(18)
            r.font.color.rgb = RGBColor(220, 220, 220)


def build_presentation():
    prs = Presentation()

    add_title_slide(prs)

    add_bulleted_slide(
        prs,
        "Project Introduction",
        [
            "AI Study Assistant is a full-stack educational web application.",
            "It helps students generate concise notes for any topic.",
            "Users can save and revisit previous notes from dashboard history.",
            "Project combines FastAPI, React, SQLAlchemy, and OpenRouter AI.",
        ],
    )

    add_bulleted_slide(
        prs,
        "Problem Statement",
        [
            "Students spend too much time creating notes manually.",
            "Generated content is often not stored in a structured way.",
            "Most generic AI tools do not provide personal note history.",
            "Need a simple system for secure, topic-wise, reusable notes.",
        ],
    )

    add_bulleted_slide(
        prs,
        "Objectives",
        [
            "Implement secure signup and login using JWT authentication.",
            "Generate AI-based notes from user-entered topics.",
            "Store notes in a database linked to each user account.",
            "Provide clean dashboard UI with notes history and quick access.",
        ],
    )

    add_two_content_slide(
        prs,
        "Technology Stack",
        [
            "Frontend: React + Vite",
            "Styling: Tailwind CSS",
            "API Calls: Axios",
            "Routing: React Router",
        ],
        [
            "Backend: FastAPI (Python)",
            "ORM: SQLAlchemy",
            "Auth: JWT + OAuth2 form login",
            "AI Provider: OpenRouter API",
        ],
    )

    add_bulleted_slide(
        prs,
        "System Architecture",
        [
            "Frontend collects input topic and sends authenticated requests.",
            "Backend validates token and identifies current user.",
            "AI service module calls OpenRouter and receives generated notes.",
            "Backend stores notes and returns response to dashboard.",
        ],
    )

    add_bulleted_slide(
        prs,
        "Database Design",
        [
            "users table: id, email, password (hashed).",
            "notes table: id, topic, content, user_id.",
            "Each note is mapped to a specific authenticated user.",
            "Design is simple, scalable, and easy to query.",
        ],
    )

    add_bulleted_slide(
        prs,
        "Core API Endpoints",
        [
            "POST /auth/signup  -> Create new account",
            "POST /auth/login  -> Return access token",
            "POST /notes/generate  -> Generate and save notes",
            "GET /notes/history  -> Fetch user note history",
        ],
    )

    add_image_slide(
        prs,
        "Screenshot: Signup Page",
        BASE_DIR / "screenshots" / "signup.png",
        "Figure: User registration page",
    )

    add_image_slide(
        prs,
        "Screenshot: Dashboard with AI Output",
        BASE_DIR / "screenshots" / "dashboard-generated.png",
        "Figure: Dashboard after generating notes on topic 'photosynthesis'",
    )

    add_bulleted_slide(
        prs,
        "Testing and Results",
        [
            "Signup and login tested with valid and invalid credentials.",
            "Protected routes validated with bearer token authentication.",
            "Note generation tested with multiple topics.",
            "History retrieval confirmed after successful note creation.",
        ],
    )

    add_bulleted_slide(
        prs,
        "Conclusion and Future Scope",
        [
            "Project successfully delivers AI-assisted personalized note generation.",
            "System is modular and suitable for academic demonstration.",
            "Future scope: PDF export, search/filter, categories, streaming output.",
            "Can be extended to a broader smart learning platform.",
        ],
    )

    prs.save(OUTPUT_FILE)
    print(OUTPUT_FILE)


if __name__ == "__main__":
    build_presentation()
